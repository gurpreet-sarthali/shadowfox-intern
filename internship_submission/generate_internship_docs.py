from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import date

# -------------------------
# Generate Word report
# -------------------------
report = Document()
report.styles['Normal'].font.name = 'Calibri'
report.styles['Normal'].font.size = Pt(11)

# Title
p = report.add_paragraph()
p.alignment = 1
run = p.add_run('Internship Report')
run.bold = True
run.font.size = Pt(24)

report.add_paragraph('')
report.add_paragraph('Prepared for Internship Documentation and Submission')
report.add_paragraph('Date: ' + date.today().strftime('%d %B %Y'))
report.add_paragraph('')

# Section 1
report.add_heading('1. Introduction', level=1)
report.add_paragraph('This internship report summarizes the work completed during the internship period. The goal of this internship was to build practical knowledge in data analysis, machine learning, and business intelligence through real datasets and hands-on practice. The work was performed using Python and relevant libraries for data processing, visualization, and model evaluation.')

# Section 2
report.add_heading('2. Internship Objectives', level=1)
report.add_paragraph('The main objectives of this internship were to:')
for item in [
    'Understand and work with real-world datasets.',
    'Apply data cleaning and preprocessing steps.',
    'Perform exploratory data analysis to extract meaningful insights.',
    'Build and evaluate machine learning models.',
    'Analyze business performance using sales and profit data.',
    'Prepare documentation and presentation material for submission.'
]:
    report.add_paragraph(item, style='List Bullet')

# Section 3 basic task
report.add_heading('3. Task 1: Boston House Price Prediction', level=1)
report.add_paragraph('The first major task involved predicting Boston house prices using the dataset available in the project folder. The dataset contained multiple economic and housing-related variables such as crime rate, number of rooms, tax, and pollution indicators. The target variable was MEDV, which represents the median value of owner-occupied homes in thousands of dollars.')
report.add_paragraph('The workflow included data inspection, missing value handling, duplicate checking, feature selection, train-test splitting, and model comparison. Several regression models were trained, including Linear Regression, Decision Tree Regressor, Random Forest Regressor, and Gradient Boosting Regressor.')
report.add_paragraph('The models were compared using MAE, MSE, RMSE, and R-squared. The best-performing model was the Gradient Boosting Regressor, which gave the strongest results on the test set.')

report.add_paragraph('Key evaluation results:')
for item in [
    'MAE: 1.9187',
    'MSE: 7.3028',
    'RMSE: 2.7024',
    'R²: 0.9004'
]:
    report.add_paragraph(item, style='List Bullet')

report.add_paragraph('This task showed that house price prediction can be effectively handled with machine learning, especially when the dataset is cleaned properly and model performance is compared systematically.')

# Section 4 intermediate task
report.add_heading('4. Task 2: Store Sales and Profit Analysis', level=1)
report.add_paragraph('The second task focused on business analytics using the Sample-Superstore dataset. The objective was to analyze sales and profit patterns across regions, product categories, customer segments, and sub-categories so that meaningful business insights could be extracted.')
report.add_paragraph('The analysis included sales and profit summaries, region-wise comparison, category-level performance, customer behavior, and monthly trends. Visualizations were used to make the results clearer and easier to interpret.')

report.add_paragraph('Key findings:')
for item in [
    'Total sales: $2,297,200.86',
    'Total profit: $286,397.02',
    'Profit margin: 12.47%',
    'West region produced the highest profit',
    'Technology was the most profitable category',
    'Consumer segment contributed the largest sales and profit',
    'There were 1,871 negative-profit transactions that need review.'
]:
    report.add_paragraph(item, style='List Bullet')

report.add_paragraph('This task was useful in understanding how data analytics can help businesses identify high-performing segments and reduce loss-making activities.')

# Section 5 tools
report.add_heading('5. Tools and Technologies Used', level=1)
report.add_paragraph('The internship tasks were completed using Python along with major libraries such as pandas, NumPy, matplotlib, seaborn, scikit-learn, and Jupyter Notebook. These tools were used for data manipulation, statistical analysis, model building, and visual reporting.')

# Section 6 learnings
report.add_heading('6. Learning and Outcome', level=1)
report.add_paragraph('This internship provided practical experience in working with real datasets, understanding business and ML problems, and applying structured problem-solving methods. It strengthened my ability to clean data, interpret results, compare models, and communicate findings in a clear and professional way. The experience also improved my understanding of how analytics can support decision-making in real-world settings.')

# Section 7 conclusion
report.add_heading('7. Conclusion', level=1)
report.add_paragraph('Overall, both tasks were successfully completed using real datasets from the project workspace. The work covered both machine learning and business intelligence aspects, which together created a strong practical learning experience. This internship work reflects the effort, technical learning, and problem-solving approach that were developed during the internship period.')

report.save(r'c:\Users\Lenovo\Desktop\my all\code\shadowfox\internship_submission\Internship_Report.docx')

# -------------------------
# Generate PPT
# -------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1 - Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Internship Report and Project Summary'
slide.placeholders[1].text = 'Data Analysis and Machine Learning Internship\nPrepared using completed project work'

# Slide 2 - Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Internship Overview'
body = slide.shapes.placeholders[1].text_frame
body.text = 'This internship involved hands-on work with real-world datasets and practical problem-solving in data analysis and machine learning.\n\nThe work included:'
body.add_paragraph()
for item in [
    'Boston house price prediction using regression models',
    'Store sales and profit analysis using business data',
    'Data cleaning, EDA, modeling, and interpretation',
    'Visualization and report preparation'
]:
    p = body.add_paragraph()
    p.text = item
    p.level = 0

# Slide 3 - Boston task summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Task 1: Boston House Price Prediction'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Dataset used: HousingData.csv\n\nKey work done:\n- Checked dataset shape, target variable, missing values, and duplicates\n- Trained multiple regression models\n- Compared model performance using MAE, MSE, RMSE, and R²\n\nBest model: Gradient Boosting Regressor\nMAE = 1.9187 | MSE = 7.3028 | RMSE = 2.7024 | R² = 0.9004'

# Slide 4 - Model comparison table style using text boxes
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Model Comparison'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Model                           MAE      MSE      RMSE      R²\nLinear Regression           3.1476   24.9834   4.9983   0.6593\nDecision Tree               3.1529   26.1347   5.1122   0.6436\nRandom Forest               2.0743    9.1011   3.0168   0.8759\nGradient Boosting           1.9187    7.3028   2.7024   0.9004'

# Slide 5 - Sales task overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Task 2: Store Sales and Profit Analysis'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Dataset used: Sample-Superstore.csv\n\nKey findings:\n- Total Sales: $2,297,200.86\n- Total Profit: $286,397.02\n- Profit margin: 12.47%\n- West region was the strongest performer\n- Technology category delivered the highest profit\n- Consumer segment generated the highest sales and profit'

# Slide 6 - Business insights
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Business Insights'
body = slide.shapes.placeholders[1].text_frame
body.text = 'The business analysis revealed important patterns for decision-making:\n\n- The West and East regions are leading contributors to revenue\n- Technology products generate the strongest profit potential\n- Consumer customers drive the largest portion of sales\n- Loss-making orders should be reviewed to reduce unnecessary costs\n- Sales and profit show a positive but moderate relationship'

# Slide 7 - Tools and outcome
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Tools, Skills and Outcome'
body = slide.shapes.placeholders[1].text_frame
body.text = 'Tools used:\n- Python\n- pandas\n- NumPy\n- matplotlib\n- seaborn\n- scikit-learn\n- Jupyter Notebook\n\nSkills developed:\n- Data cleaning and preprocessing\n- Statistical analysis and visualization\n- Regression modeling\n- Business insight generation\n- Documentation and presentation preparation'

# Slide 8 - Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = 'Conclusion'
body = slide.shapes.placeholders[1].text_frame
body.text = 'The internship work was completed successfully using real datasets and practical analytical methods. Both a machine learning project and a business analysis project were completed, and the results were verified. This work reflects the learning, effort, and problem-solving ability developed during the internship period.\n\nPrepared for internship documentation and submission.'

prs.save(r'c:\Users\Lenovo\Desktop\my all\code\shadowfox\internship_submission\Internship_Presentation.pptx')

print('Report and presentation generated successfully.')
print('Files:')
print(r'c:\Users\Lenovo\Desktop\my all\code\shadowfox\internship_submission\Internship_Report.docx')
print(r'c:\Users\Lenovo\Desktop\my all\code\shadowfox\internship_submission\Internship_Presentation.pptx')
